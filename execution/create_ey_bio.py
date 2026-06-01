"""
Create EY Bio PowerPoint for Ajinkya Pawale
Uses Alexandre Truong bio as template for formatting
Preserves exact layout, fonts, spacing, and paragraph structure
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from copy import deepcopy
import os

# Paths
TEMPLATE_PATH = "/Users/ajinkyapawale/Claude Code Projects/data/CV - Alexandre Truong 3.pptx"
OUTPUT_PATH = "/Users/ajinkyapawale/Claude Code Projects/data/Ajinkya_Pawale_EY_Bio.pptx"
PHOTO_PATH = "/Users/ajinkyapawale/Downloads/Ajinkya_Profile_Picture_Copy.jpg"

# Content for Ajinkya Pawale

# Header content (5 lines to match template structure)
HEADER_LINES = [
    ("Ajinkya Pawale", True),      # (text, bold)
    ("Senior Consultant", False),
    ("AI & Data", False),
    ("", False),                    # Empty line for spacing (matches template)
    ("Tel : 763-306-4790", False),
    ("Email\tajinkya.pawale@ey.com", False),
]

# Background content - 4 paragraphs matching template structure
BACKGROUND_PARAGRAPHS = [
    (1, "Ajinkya joined the AI & Data Practice of EY Charlotte in January 2026. He brings extensive experience in AI/ML engineering, GenAI solutions, and data pipeline development. Prior to EY, he worked as an AI Engineer at Elevance Health and as a Data Engineer at Course5 Intelligence."),
    (1, "Ajinkya holds the following qualifications:"),
    (2, "Master of Science in Data Science, Indiana University Bloomington"),
    (2, "Bachelor of Science in Computer Science, Mumbai University"),
]

# Skills content - 8 paragraphs matching template structure
SKILLS_PARAGRAPHS = [
    (1, "Healthcare, E-commerce, Retail Banking"),
    (1, "Machine Learning, Generative AI, LLM Fine-Tuning, Predictive Analytics, Deep Learning, NLP"),
    (1, "Data Engineering, Cloud Architecture, MLOps"),
    (1, "Programming and tools:"),
    (2, "Python, PySpark, R, SQL, PyTorch, TensorFlow"),
    (2, "Langchain, LlamaIndex, OpenAI and Anthropic APIs, Hugging Face"),
    (2, "AWS (S3, EC2, Lambda, EMR), Kubernetes, Docker, Airflow"),
    (1, "Databases: PostgreSQL, MongoDB, Hadoop, Hive, Qdrant"),
]

# Experience Middle Column - 9 paragraphs matching template
# All at level 0 (matching template), headers have bullets removed via XML
EXPERIENCE_MIDDLE_PARAGRAPHS = [
    (0, "Generative AI", True),  # (level, text, is_category_header)
    (0, "Led end-to-end GenAI solution for automated call summarization using OpenAI GPT models, architecting Langchain-based orchestration that achieved 55% faster response times across 350+ agents", False),
    (0, "Augmented call wrap-up API with RAG pipelines and autonomous agent workflows for validation, achieving 85% higher summary quality scores and cutting post-call correction time by two-thirds", False),
    (0, "Fine-tuned open-source LLMs (Llama) using LoRA on healthcare call transcript corpora for domain-adapted summarization, building evaluation pipelines to benchmark summarization accuracy and clinical terminology coverage against baseline GPT outputs", False),
    (0, "Architected multi-agent GPT-4.1-based document compliance system for 50+ L&D team members, reducing manual review time by 80% for large-scale training materials", False),
    (0, "Machine Learning & Predictive Analytics", True),
    (0, "Enhanced data retrieval for Inpatient Claims using PySpark, achieving 75% reduction in extraction time; enriched ML model inputs with 20+ derived features using PyTorch", False),
    (0, "Developed neural network and Bayesian optimization (BOHB) models using TensorFlow and PyTorch, boosting model efficiency by 12% and enabling early identification of at-risk students", False),
    (0, "Conducted statistical analysis in R, pandas, scikit-learn, demonstrating re-identification risks with 20% improvement in predictive accuracy", False),
]

# Experience Right Column - up to 12 paragraphs
# All at level 0 (matching template), headers have bullets removed via XML
EXPERIENCE_RIGHT_PARAGRAPHS = [
    (0, "Data Engineering & Optimization", True),
    (0, "Spearheaded automation of data pipelines in Airflow, integrating ETL, modeling, and post-processing, saving 4 hours of manual work per production run", False),
    (0, "Designed Kinesis Consumer Application for live call transcript processing with Redis cache, enhancing data retrieval speed by 40%", False),
    (0, "Migrated project data from on-premise to AWS EMR cluster, improving processing speed, scalability, and reliability", False),
    (0, "Implemented data validation tool detecting production data drift, saving nearly two full days of workload monthly", False),
    (0, "Analytics & Monitoring", True),
    (0, "Established 20+ crucial alerts and dashboards in Datadog and Splunk to monitor system health, accuracy, and performance", False),
    (0, "Performed statistical analysis on 100,000+ data entries across six sources; utilized Tableau visualizations reducing manual analysis time by 50%", False),
    (0, "Created analytics framework using pymongo to query data lake, generating insights across 250+ product categories", False),
    (0, "Testing & Quality Assurance", True),
    (0, "Orchestrated pytest framework covering 50+ test scenarios, reducing manual testing time by 3 hours per week", False),
    (0, "Executed load and stress tests using LoadRunner, coordinated API-based regression testing across feature cycles", False),
]


def update_shape_text(shape, paragraphs_data, font_name="EYInterstate Light", use_justify=True):
    """
    Update shape text while preserving formatting.
    paragraphs_data: list of (level, text) or (level, text, is_header) tuples
    """
    if not hasattr(shape, 'text_frame'):
        return

    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)

    for i, para_data in enumerate(paragraphs_data):
        if len(para_data) == 2:
            level, text = para_data
            is_header = False
        else:
            level, text, is_header = para_data

        if i < len(existing_paras):
            para = existing_paras[i]
            # Clear existing runs
            for run in list(para.runs):
                run.text = ""

            # Set paragraph level
            para.level = level

            # Add new run with text
            if para.runs:
                run = para.runs[0]
                run.text = text
            else:
                run = para.add_run()
                run.text = text

            # Set font properties
            run.font.name = font_name
            run.font.size = Pt(10)
            run.font.bold = is_header if len(para_data) == 3 else None

            # Set alignment if needed
            if use_justify:
                para.alignment = PP_ALIGN.JUSTIFY


def update_header(shape, lines):
    """Update header text box with name, title, etc."""
    if not hasattr(shape, 'text_frame'):
        return

    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)

    for i, (text, is_bold) in enumerate(lines):
        if i < len(existing_paras):
            para = existing_paras[i]
            # Clear existing runs
            for run in list(para.runs):
                run.text = ""

            # Add new run
            if para.runs:
                run = para.runs[0]
                run.text = text
            else:
                run = para.add_run()
                run.text = text

            run.font.name = "EYInterstate"
            run.font.size = Pt(9)
            run.font.bold = is_bold
            run.font.underline = False


def update_background(shape, paragraphs_data):
    """Update background section with proper level formatting."""
    if not hasattr(shape, 'text_frame'):
        return

    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)

    for i, (level, text) in enumerate(paragraphs_data):
        if i < len(existing_paras):
            para = existing_paras[i]
            # Clear existing runs
            for run in list(para.runs):
                run.text = ""

            para.level = level

            if para.runs:
                run = para.runs[0]
                run.text = text
            else:
                run = para.add_run()
                run.text = text

            run.font.name = "EYInterstate"
            run.font.size = Pt(10)
            run.font.underline = False
            run.font.bold = False

            # Set consistent spacing for education qualifications (indices 2 and 3)
            if i >= 2:
                para.space_before = Pt(0)
                para.space_after = Pt(0)

            para.alignment = PP_ALIGN.JUSTIFY


def update_skills(shape, paragraphs_data):
    """Update skills section with proper level formatting."""
    if not hasattr(shape, 'text_frame'):
        return

    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)

    for i, (level, text) in enumerate(paragraphs_data):
        if i < len(existing_paras):
            para = existing_paras[i]
            # Clear existing runs
            for run in list(para.runs):
                run.text = ""

            para.level = level

            if para.runs:
                run = para.runs[0]
                run.text = text
            else:
                run = para.add_run()
                run.text = text

            run.font.name = "EYInterstate"
            run.font.size = Pt(10)
            run.font.underline = False
            run.font.bold = False


def remove_bullet(para):
    """Remove bullet from a paragraph (for headers) by removing buChar."""
    pPr = para._p.get_or_add_pPr()
    # Remove buChar element if present (this removes the bullet)
    for child in list(pPr):
        if 'buChar' in child.tag or 'buAutoNum' in child.tag:
            pPr.remove(child)


def ensure_bullet(para):
    """Ensure paragraph has a bullet character (for bullet points)."""
    pPr = para._p.get_or_add_pPr()
    # Check if buChar already exists
    has_buChar = any('buChar' in child.tag for child in pPr)
    if not has_buChar:
        # Add bullet character (dash style like template)
        buChar = parse_xml('<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="–"/>')
        pPr.append(buChar)


def update_experience(shape, paragraphs_data):
    """Update experience section with category headers (underlined, no bullets) and bullet points."""
    if not hasattr(shape, 'text_frame'):
        return

    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)

    for i, para_data in enumerate(paragraphs_data):
        level, text, is_header = para_data

        if i < len(existing_paras):
            para = existing_paras[i]
            # Clear existing runs
            for run in list(para.runs):
                run.text = ""

            para.level = level

            if para.runs:
                run = para.runs[0]
                run.text = text
            else:
                run = para.add_run()
                run.text = text

            run.font.name = "EYInterstate Light"
            run.font.size = Pt(10)

            # Headers: underlined, no bullet; Bullet points: not underlined, with bullet
            if is_header:
                run.font.underline = True
                run.font.bold = False
                remove_bullet(para)
            else:
                run.font.underline = False
                run.font.bold = False
                ensure_bullet(para)

            para.alignment = PP_ALIGN.JUSTIFY


def update_photo(slide, photo_path):
    """Replace existing photo or add new one at correct position."""
    # Template photo position and size
    left = Emu(634910)   # 0.694"
    top = Emu(222867)    # 0.244"
    width = Emu(777600)  # 0.85"
    height = Emu(1036800)  # 1.134"

    # Remove existing picture
    pictures_to_remove = []
    for shape in slide.shapes:
        if 'picture' in shape.name.lower():
            pictures_to_remove.append(shape)

    for pic in pictures_to_remove:
        sp = pic._element
        sp.getparent().remove(sp)

    # Add new picture
    if os.path.exists(photo_path):
        slide.shapes.add_picture(photo_path, left, top, width, height)
        return True
    return False


# Main execution
print("Loading template...")
prs = Presentation(TEMPLATE_PATH)
slide = prs.slides[0]

print("Processing shapes...")
for shape in slide.shapes:
    shape_name = shape.name

    if shape_name == "Text Box 7":
        print(f"  Updating header: {shape_name}")
        update_header(shape, HEADER_LINES)

    elif shape_name == "Rectangle 18":
        print(f"  Updating background: {shape_name}")
        update_background(shape, BACKGROUND_PARAGRAPHS)

    elif shape_name == "Rectangle 21":
        print(f"  Updating skills: {shape_name}")
        update_skills(shape, SKILLS_PARAGRAPHS)

    elif shape_name == "Rectangle 15":
        print(f"  Updating experience (middle): {shape_name}")
        update_experience(shape, EXPERIENCE_MIDDLE_PARAGRAPHS)

    elif shape_name == "Rectangle 28":
        print(f"  Updating experience (right): {shape_name}")
        update_experience(shape, EXPERIENCE_RIGHT_PARAGRAPHS)

print("Updating photo...")
if update_photo(slide, PHOTO_PATH):
    print("  Photo added successfully")
else:
    print(f"  Photo not found at: {PHOTO_PATH}")

print("Saving presentation...")
prs.save(OUTPUT_PATH)
print(f"\nBio saved to: {OUTPUT_PATH}")
